#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import re
import shutil
import subprocess
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional
import xml.etree.ElementTree as ET

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import docx
except Exception:
    docx = None


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".ppt",
    ".pptx",
    ".doc",
    ".docx",
    ".xls",
    ".xlsx",
    ".md",
}

BUSINESS_RULES = {
    "finished_inspection_process": [
        "finished product inspection process",
        "成品查貨流程",
        "成品查货流程",
        "3.3a",
        "3.3b",
    ],
    "semi_finished_inspection_process": [
        "semi-finished product inspection process",
        "半成品查貨流程",
        "半成品查货流程",
        "4.7",
    ],
    "quality_standard": [
        "quality standards",
        "品質標準",
        "品质标准",
        "抽樣",
        "抽样",
        "aql",
    ],
    "defect_catalog": [
        "defect analysis table",
        "疵點",
        "疵点",
        "缺陷",
        "破洞",
        "污渍",
        "污漬",
        "色差",
        "跳线",
        "跳線",
    ],
    "inspection_guideline": [
        "product inspection and inventory check guidelines",
        "產品檢查及查貨工作指引",
        "产品检查及查货工作指引",
        "工作指引",
    ],
    "work_guideline": [
        "work guidelines",
        "品質保證部工作指引",
        "品质保证部工作指引",
    ],
    "inspection_plan": [
        "order inspection plan development",
        "制定订单查货计划",
        "制定訂單查貨計劃",
        "inspection plan",
    ],
}

CATEGORY_PATH_HINTS = {
    "finished_inspection_process": [
        "finished product inspection process",
    ],
    "semi_finished_inspection_process": [
        "semi-finished product inspection process",
    ],
    "quality_standard": [
        "quality standards",
        "品质标准",
        "品質標準",
    ],
    "defect_catalog": [
        "defect analysis table",
    ],
    "inspection_guideline": [
        "product inspection and inventory check guidelines",
    ],
    "work_guideline": [
        "work guidelines",
    ],
    "inspection_plan": [
        "order inspection plan development",
    ],
}

CATEGORY_TO_FOLDER = {
    "finished_inspection_process": "Finished Product Inspection Process",
    "semi_finished_inspection_process": "Semi-finished Product Inspection Process",
    "quality_standard": "Quality Standards",
    "defect_catalog": "Defect Analysis Table",
    "inspection_guideline": "Product Inspection and Inventory Check Guidelines",
    "work_guideline": "Work guidelines",
    "inspection_plan": "Order inspection plan development",
    "uncategorized": "Uncategorized",
}

LOCATION_LABELS = {
    "pdf": "page",
    "ppt": "slide",
    "pptx": "slide",
    "word": "section",
    "excel": "sheet",
    "markdown": "section",
}

NS_MAIN = {"main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
NS_REL = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}


@dataclass
class ExtractionResult:
    text: str
    sections: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    extractor: str
    status: str
    extracted_from: Optional[str] = None
    requires_ocr: bool = False
    notes: Optional[str] = None


def utc_now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def normalize_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def stable_doc_id(source_identity: str, tech_type: str) -> str:
    digest = hashlib.sha1(source_identity.encode("utf-8")).hexdigest()[:12]
    return f"{tech_type}-{digest}"


def detect_tech_type(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".ppt", ".pptx"}:
        return "pptx" if suffix == ".pptx" else "ppt"
    if suffix in {".doc", ".docx"}:
        return "word"
    if suffix in {".xls", ".xlsx"}:
        return "excel"
    if suffix == ".pdf":
        return "pdf"
    return "markdown"


def classify_business(path: Path, sample_text: str = "") -> Dict[str, Any]:
    path_parts = [part.strip().lower() for part in path.parts]
    path_text = " ".join(path.parts)
    for category, hints in CATEGORY_PATH_HINTS.items():
        matched_path_hints = [hint for hint in hints if hint.lower() in path_parts]
        if matched_path_hints:
            return {
                "business_category": category,
                "matched_keywords": matched_path_hints,
                "target_folder": CATEGORY_TO_FOLDER[category],
            }

    haystack = path_text + "\n" + sample_text[:4000]
    lowered = haystack.lower()
    scores: List[tuple[int, str, List[str]]] = []
    for category, keywords in BUSINESS_RULES.items():
        hits = [keyword for keyword in keywords if keyword.lower() in lowered]
        scores.append((len(hits), category, hits))
    scores.sort(reverse=True)
    top_score, top_category, hits = scores[0]
    if top_score == 0:
        top_category = "uncategorized"
        hits = []
    return {
        "business_category": top_category,
        "matched_keywords": hits,
        "target_folder": CATEGORY_TO_FOLDER[top_category],
    }


def relative_to(base: Path, target: Path) -> str:
    try:
        return str(target.relative_to(base))
    except ValueError:
        return str(target.resolve())


def ensure_directory(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def unique_destination(path: Path) -> Path:
    if not path.exists():
        return path
    stem = path.stem
    suffix = path.suffix
    counter = 1
    while True:
        candidate = path.with_name(f"{stem}__{counter}{suffix}")
        if not candidate.exists():
            return candidate
        counter += 1


def maybe_copy_into_library(skill_root: Path, source_path: Path, category: str) -> tuple[Path, bool]:
    knowledge_base = skill_root / "references" / "knowledge_base"
    try:
        source_path.relative_to(knowledge_base)
        return source_path, False
    except ValueError:
        pass

    target_root = knowledge_base / CATEGORY_TO_FOLDER.get(category, "Uncategorized") / "auto_ingested"
    ensure_directory(target_root)
    destination = unique_destination(target_root / source_path.name)
    shutil.copy2(source_path, destination)
    return destination, True


def relocate_library_copy(skill_root: Path, library_path: Path, category: str) -> Path:
    knowledge_base = skill_root / "references" / "knowledge_base"
    target_root = knowledge_base / CATEGORY_TO_FOLDER.get(category, "Uncategorized") / "auto_ingested"
    ensure_directory(target_root)
    destination = unique_destination(target_root / library_path.name)
    if destination == library_path:
        return library_path
    shutil.move(str(library_path), str(destination))
    return destination


def load_registry(path: Path) -> Dict[str, Dict[str, Any]]:
    records: Dict[str, Dict[str, Any]] = {}
    if not path.exists():
        return records
    with path.open("r", encoding="utf-8") as handle:
        for line in handle:
            line = line.strip()
            if not line:
                continue
            record = json.loads(line)
            records[record["source_identity"]] = record
    return records


def save_registry(path: Path, records: Dict[str, Dict[str, Any]]) -> None:
    ensure_directory(path.parent)
    ordered = [records[key] for key in sorted(records)]
    with path.open("w", encoding="utf-8") as handle:
        for record in ordered:
            handle.write(json.dumps(record, ensure_ascii=False) + "\n")


def chunk_text(text: str, max_chars: int = 1200) -> List[str]:
    paragraphs = [paragraph.strip() for paragraph in text.split("\n\n") if paragraph.strip()]
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0
    for paragraph in paragraphs:
        addition = len(paragraph) + (2 if current else 0)
        if current and current_len + addition > max_chars:
            chunks.append("\n\n".join(current))
            current = [paragraph]
            current_len = len(paragraph)
            continue
        current.append(paragraph)
        current_len += addition
    if current:
        chunks.append("\n\n".join(current))
    return chunks or ([text] if text else [])


def build_chunks(
    doc_id: str,
    source_rel_path: str,
    tech_type: str,
    business_category: str,
    sections: List[Dict[str, Any]],
    fallback_text: str,
) -> List[Dict[str, Any]]:
    chunks: List[Dict[str, Any]] = []
    location_label = LOCATION_LABELS.get(tech_type, "section")
    if sections:
        for index, section in enumerate(sections, start=1):
            body = normalize_text(section.get("body", ""))
            if not body:
                continue
            heading = section.get("heading") or f"{location_label.title()} {index}"
            location_value = section.get("location") or index
            chunks.append(
                {
                    "chunk_id": f"{doc_id}:{index:04d}",
                    "doc_id": doc_id,
                    "source_path": source_rel_path,
                    "file_type": tech_type,
                    "business_category": business_category,
                    "heading": heading,
                    "location_label": location_label,
                    "location_value": location_value,
                    "text": body,
                }
            )
        if chunks:
            return chunks
    for index, piece in enumerate(chunk_text(fallback_text), start=1):
        chunks.append(
            {
                "chunk_id": f"{doc_id}:{index:04d}",
                "doc_id": doc_id,
                "source_path": source_rel_path,
                "file_type": tech_type,
                "business_category": business_category,
                "heading": f"{location_label.title()} {index}",
                "location_label": location_label,
                "location_value": index,
                "text": normalize_text(piece),
            }
        )
    return chunks


def markdown_table(rows: List[List[str]]) -> str:
    cleaned = [[cell.strip() for cell in row] for row in rows if any(cell.strip() for cell in row)]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    padded = [row + [""] * (width - len(row)) for row in cleaned]
    header = padded[0]
    divider = ["---"] * width
    body = padded[1:] or [[""] * width]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(divider) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def extract_markdown(path: Path) -> ExtractionResult:
    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    sections: List[Dict[str, Any]] = []
    current_heading = "Document"
    current_lines: List[str] = []
    for line in raw_text.splitlines():
        if line.startswith("#"):
            if current_lines:
                sections.append(
                    {
                        "heading": current_heading,
                        "location": current_heading,
                        "body": "\n".join(current_lines).strip(),
                    }
                )
                current_lines = []
            current_heading = line.lstrip("#").strip() or "Document"
            continue
        current_lines.append(line)
    if current_lines:
        sections.append(
            {
                "heading": current_heading,
                "location": current_heading,
                "body": "\n".join(current_lines).strip(),
            }
        )
    return ExtractionResult(
        text=normalize_text(raw_text),
        sections=sections,
        tables=[],
        extractor="markdown-direct",
        status="ok",
    )


def extract_pdf(path: Path) -> ExtractionResult:
    sections: List[Dict[str, Any]] = []
    pages_with_text = 0
    if pdfplumber is not None:
        with pdfplumber.open(str(path)) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = normalize_text(page.extract_text() or "")
                if text:
                    pages_with_text += 1
                sections.append(
                    {
                        "heading": f"Page {index}",
                        "location": index,
                        "body": text,
                    }
                )
    elif PdfReader is not None:
        reader = PdfReader(str(path))
        for index, page in enumerate(reader.pages, start=1):
            text = normalize_text(page.extract_text() or "")
            if text:
                pages_with_text += 1
            sections.append(
                {
                    "heading": f"Page {index}",
                    "location": index,
                    "body": text,
                }
            )
    else:
        return ExtractionResult(
            text="",
            sections=[],
            tables=[],
            extractor="pdf-unavailable",
            status="extractor_missing",
            notes="Neither pdfplumber nor pypdf is available.",
        )

    full_text = normalize_text("\n\n".join(section["body"] for section in sections if section["body"]))
    requires_ocr = not bool(full_text)
    return ExtractionResult(
        text=full_text,
        sections=sections,
        tables=[],
        extractor="pdfplumber" if pdfplumber is not None else "pypdf",
        status="ok" if full_text else "needs_ocr",
        requires_ocr=requires_ocr,
        notes="No text extracted from PDF pages." if requires_ocr else None,
    )


def extract_pptx(path: Path) -> ExtractionResult:
    if Presentation is None:
        return ExtractionResult(
            text="",
            sections=[],
            tables=[],
            extractor="pptx-unavailable",
            status="extractor_missing",
            notes="python-pptx is not available.",
        )

    presentation = Presentation(str(path))
    sections: List[Dict[str, Any]] = []
    tables: List[Dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = ""
        text_runs: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text = normalize_text(shape.text or "")
                if not text:
                    continue
                if not title:
                    title = text.splitlines()[0]
                text_runs.append(text)
            if hasattr(shape, "has_table") and shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([normalize_text(cell.text or "") for cell in row.cells])
                table_md = markdown_table(rows)
                if table_md:
                    text_runs.append(table_md)
                    tables.append(
                        {
                            "name": f"slide_{index}_table_{len(tables) + 1}",
                            "rows": rows,
                        }
                    )
        notes_text = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                notes_text = normalize_text(slide.notes_slide.notes_text_frame.text or "")
            except Exception:
                notes_text = ""
        body_parts = []
        if title:
            body_parts.append(f"# {title}")
        body_parts.extend(text_runs)
        if notes_text:
            body_parts.append("Notes\n" + notes_text)
        sections.append(
            {
                "heading": title or f"Slide {index}",
                "location": index,
                "body": normalize_text("\n\n".join(part for part in body_parts if part)),
            }
        )

    full_text = normalize_text("\n\n".join(section["body"] for section in sections if section["body"]))
    status = "ok" if full_text else "needs_ocr"
    return ExtractionResult(
        text=full_text,
        sections=sections,
        tables=tables,
        extractor="python-pptx",
        status=status,
        requires_ocr=not bool(full_text),
        notes="Slides contain no extractable text." if not full_text else None,
    )


def extract_ppt(path: Path) -> ExtractionResult:
    pdf_fallback = path.with_suffix(".pdf")
    if pdf_fallback.exists():
        result = extract_pdf(pdf_fallback)
        result.extractor = "paired-pdf-for-ppt"
        result.extracted_from = str(pdf_fallback)
        return result
    return ExtractionResult(
        text="",
        sections=[],
        tables=[],
        extractor="ppt-metadata-only",
        status="manual_review",
        requires_ocr=True,
        notes="Legacy .ppt needs a paired PDF or manual conversion before reliable extraction.",
    )


def extract_docx(path: Path) -> ExtractionResult:
    if docx is None:
        return ExtractionResult(
            text="",
            sections=[],
            tables=[],
            extractor="docx-unavailable",
            status="extractor_missing",
            notes="python-docx is not available.",
        )
    document = docx.Document(str(path))
    sections: List[Dict[str, Any]] = []
    tables: List[Dict[str, Any]] = []
    current_heading = path.stem
    current_lines: List[str] = []

    def flush_section() -> None:
        nonlocal current_lines
        body = normalize_text("\n".join(current_lines))
        if body:
            sections.append(
                {
                    "heading": current_heading,
                    "location": current_heading,
                    "body": body,
                }
            )
        current_lines = []

    for paragraph in document.paragraphs:
        text = normalize_text(paragraph.text or "")
        if not text:
            continue
        style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style_name.startswith("heading"):
            flush_section()
            current_heading = text
            continue
        current_lines.append(text)

    for index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            rows.append([normalize_text(cell.text or "") for cell in row.cells])
        table_md = markdown_table(rows)
        if table_md:
            flush_section()
            sections.append(
                {
                    "heading": f"Table {index}",
                    "location": f"table-{index}",
                    "body": table_md,
                }
            )
            tables.append({"name": f"table_{index}", "rows": rows})

    flush_section()
    full_text = normalize_text("\n\n".join(section["body"] for section in sections if section["body"]))
    return ExtractionResult(
        text=full_text,
        sections=sections,
        tables=tables,
        extractor="python-docx",
        status="ok" if full_text else "empty",
    )


def extract_doc(path: Path) -> ExtractionResult:
    try:
        result = subprocess.run(
            ["/usr/bin/textutil", "-convert", "txt", "-stdout", str(path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        text = normalize_text(result.stdout.decode("utf-8", errors="ignore"))
    except Exception as exc:
        return ExtractionResult(
            text="",
            sections=[],
            tables=[],
            extractor="textutil-doc",
            status="manual_review",
            notes=f"Legacy .doc extraction failed: {exc}",
        )
    return ExtractionResult(
        text=text,
        sections=[{"heading": path.stem, "location": path.stem, "body": text}] if text else [],
        tables=[],
        extractor="textutil-doc",
        status="ok" if text else "empty",
    )


def shared_strings_from_xlsx(archive: zipfile.ZipFile) -> List[str]:
    if "xl/sharedStrings.xml" not in archive.namelist():
        return []
    root = ET.fromstring(archive.read("xl/sharedStrings.xml"))
    strings: List[str] = []
    for si in root.findall("main:si", NS_MAIN):
        text_parts = [node.text or "" for node in si.findall(".//main:t", NS_MAIN)]
        strings.append("".join(text_parts))
    return strings


def workbook_sheet_targets(archive: zipfile.ZipFile) -> List[Dict[str, str]]:
    workbook = ET.fromstring(archive.read("xl/workbook.xml"))
    rels = ET.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
    rel_map = {
        rel.attrib.get("Id"): rel.attrib.get("Target")
        for rel in rels.findall("rel:Relationship", NS_REL)
    }
    sheets = []
    for sheet in workbook.findall("main:sheets/main:sheet", NS_MAIN):
        rel_id = sheet.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        target = rel_map.get(rel_id, "")
        if not target.startswith("worksheets/"):
            continue
        sheets.append(
            {
                "name": sheet.attrib.get("name", "Sheet"),
                "target": f"xl/{target}",
            }
        )
    return sheets


def column_name(index: int) -> str:
    name = ""
    current = index
    while current:
        current, remainder = divmod(current - 1, 26)
        name = chr(65 + remainder) + name
    return name


def parse_xlsx_sheet(path: str, archive: zipfile.ZipFile, shared_strings: List[str]) -> List[List[str]]:
    root = ET.fromstring(archive.read(path))
    rows: List[List[str]] = []
    for row in root.findall(".//main:sheetData/main:row", NS_MAIN):
        cells: Dict[int, str] = {}
        max_index = 0
        for cell in row.findall("main:c", NS_MAIN):
            ref = cell.attrib.get("r", "")
            match = re.match(r"([A-Z]+)", ref)
            if not match:
                continue
            letters = match.group(1)
            index = 0
            for char in letters:
                index = index * 26 + (ord(char) - 64)
            max_index = max(max_index, index)
            value = ""
            value_node = cell.find("main:v", NS_MAIN)
            if value_node is None:
                inline_text = cell.find(".//main:t", NS_MAIN)
                value = inline_text.text if inline_text is not None and inline_text.text else ""
            else:
                raw = value_node.text or ""
                if cell.attrib.get("t") == "s":
                    try:
                        value = shared_strings[int(raw)]
                    except Exception:
                        value = raw
                else:
                    value = raw
            cells[index] = normalize_text(value)
        row_values = [cells.get(i, "") for i in range(1, max_index + 1)]
        if any(cell for cell in row_values):
            rows.append(row_values)
    return rows


def extract_xlsx(path: Path) -> ExtractionResult:
    sections: List[Dict[str, Any]] = []
    tables: List[Dict[str, Any]] = []
    with zipfile.ZipFile(path) as archive:
        shared_strings = shared_strings_from_xlsx(archive)
        for sheet in workbook_sheet_targets(archive):
            rows = parse_xlsx_sheet(sheet["target"], archive, shared_strings)
            if not rows:
                continue
            preview_rows = rows[:50]
            header = preview_rows[0] if preview_rows else [column_name(i + 1) for i in range(len(preview_rows[0]))]
            summary_lines = [f"Sheet: {sheet['name']}", f"Rows extracted: {len(rows)}", ""]
            summary_lines.append(markdown_table(preview_rows if preview_rows else [header]))
            sections.append(
                {
                    "heading": sheet["name"],
                    "location": sheet["name"],
                    "body": normalize_text("\n".join(summary_lines)),
                }
            )
            tables.append({"name": sheet["name"], "rows": rows})
    full_text = normalize_text("\n\n".join(section["body"] for section in sections if section["body"]))
    return ExtractionResult(
        text=full_text,
        sections=sections,
        tables=tables,
        extractor="xlsx-xml",
        status="ok" if full_text else "empty",
    )


def extract_xls(path: Path) -> ExtractionResult:
    return ExtractionResult(
        text="",
        sections=[],
        tables=[],
        extractor="xls-metadata-only",
        status="manual_review",
        notes="Legacy .xls is not yet parsed. Convert to .xlsx for structured ingestion.",
    )


def extract_file(path: Path) -> ExtractionResult:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".ppt":
        return extract_ppt(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".doc":
        return extract_doc(path)
    if suffix == ".xlsx":
        return extract_xlsx(path)
    if suffix == ".xls":
        return extract_xls(path)
    return extract_markdown(path)


def write_table_csv(path: Path, rows: List[List[str]]) -> None:
    ensure_directory(path.parent)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerows(rows)


def write_processed_outputs(
    skill_root: Path,
    doc_id: str,
    source_rel_path: str,
    tech_type: str,
    category_info: Dict[str, Any],
    extraction: ExtractionResult,
    file_hash: str,
) -> Dict[str, Any]:
    processed_root = skill_root / "references" / "processed" / doc_id
    tables_root = processed_root / "tables"
    assets_root = processed_root / "assets"
    ensure_directory(tables_root)
    ensure_directory(assets_root)

    chunks = build_chunks(
        doc_id=doc_id,
        source_rel_path=source_rel_path,
        tech_type=tech_type,
        business_category=category_info["business_category"],
        sections=extraction.sections,
        fallback_text=extraction.text,
    )

    for index, table in enumerate(extraction.tables, start=1):
        table_name = table.get("name") or f"table_{index}"
        safe_name = re.sub(r"[^A-Za-z0-9._-]+", "_", table_name).strip("_") or f"table_{index}"
        write_table_csv(tables_root / f"{safe_name}.csv", table.get("rows", []))

    metadata = {
        "doc_id": doc_id,
        "source_path": source_rel_path,
        "file_type": tech_type,
        "business_category": category_info["business_category"],
        "matched_keywords": category_info["matched_keywords"],
        "sha256": file_hash,
        "extractor": extraction.extractor,
        "status": extraction.status,
        "requires_ocr": extraction.requires_ocr,
        "extracted_from": extraction.extracted_from,
        "notes": extraction.notes,
        "chunk_count": len(chunks),
        "updated_at": utc_now_iso(),
    }

    ensure_directory(processed_root)
    (processed_root / "metadata.json").write_text(
        json.dumps(metadata, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    (processed_root / "text.md").write_text(extraction.text + ("\n" if extraction.text else ""), encoding="utf-8")
    with (processed_root / "chunks.jsonl").open("w", encoding="utf-8") as handle:
        for chunk in chunks:
            handle.write(json.dumps(chunk, ensure_ascii=False) + "\n")
    return {
        "processed_root": relative_to(skill_root, processed_root),
        "metadata": metadata,
        "chunks": chunks,
    }


def discover_files(path: Path) -> Iterable[Path]:
    if path.is_file():
        if not path.name.startswith(".") and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path
        return
    for candidate in sorted(path.rglob("*")):
        if any(part.startswith(".") for part in candidate.parts):
            continue
        if "archive" in candidate.parts:
            continue
        if candidate.is_file() and candidate.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield candidate


def ingest_path(skill_root: Path, source: Path, force: bool) -> List[Dict[str, Any]]:
    registry_path = skill_root / "references" / "registry.jsonl"
    registry = load_registry(registry_path)
    ingested_records: List[Dict[str, Any]] = []

    for candidate in discover_files(source):
        initial_category = classify_business(candidate)
        library_path, copied_into_library = maybe_copy_into_library(
            skill_root,
            candidate,
            initial_category["business_category"],
        )
        source_identity = relative_to(skill_root, library_path)
        tech_type = detect_tech_type(library_path)
        file_hash = sha256_file(library_path)
        existing = registry.get(source_identity)
        if existing and existing.get("sha256") == file_hash and not force:
            ingested_records.append(existing)
            continue

        extraction = extract_file(library_path)
        refined_category = classify_business(library_path, extraction.text)
        if copied_into_library and refined_category["business_category"] != initial_category["business_category"]:
            library_path = relocate_library_copy(
                skill_root,
                library_path,
                refined_category["business_category"],
            )
            source_identity = relative_to(skill_root, library_path)
        doc_id = stable_doc_id(source_identity, tech_type)
        processed = write_processed_outputs(
            skill_root=skill_root,
            doc_id=doc_id,
            source_rel_path=source_identity,
            tech_type=tech_type,
            category_info=refined_category,
            extraction=extraction,
            file_hash=file_hash,
        )

        record = {
            "doc_id": doc_id,
            "source_identity": source_identity,
            "source_path": source_identity,
            "origin_path": relative_to(skill_root, candidate.resolve()),
            "file_name": library_path.name,
            "sha256": file_hash,
            "file_type": tech_type,
            "business_category": refined_category["business_category"],
            "matched_keywords": refined_category["matched_keywords"],
            "status": extraction.status,
            "requires_ocr": extraction.requires_ocr,
            "extractor": extraction.extractor,
            "notes": extraction.notes,
            "processed_root": processed["processed_root"],
            "chunk_count": len(processed["chunks"]),
            "updated_at": utc_now_iso(),
            "extracted_from": extraction.extracted_from,
        }
        registry[source_identity] = record
        ingested_records.append(record)

    save_registry(registry_path, registry)
    return ingested_records


def main() -> int:
    parser = argparse.ArgumentParser(description="Ingest and preprocess skill knowledge files.")
    parser.add_argument(
        "--source",
        default=None,
        help="File or directory to ingest. Defaults to the skill inbox directory.",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Rebuild processed outputs even if the file hash did not change.",
    )
    parser.add_argument(
        "--rebuild-index",
        action="store_true",
        help="Regenerate references/knowledge_index.md after ingestion.",
    )
    args = parser.parse_args()

    skill_root = Path(__file__).resolve().parents[1]
    default_source = skill_root / "inbox"
    source = Path(args.source).resolve() if args.source else default_source.resolve()

    ensure_directory(skill_root / "references" / "processed")
    ensure_directory(skill_root / "inbox")
    ensure_directory(skill_root / "references")

    records = ingest_path(skill_root, source, force=args.force)
    if args.rebuild_index:
        subprocess.run(
            ["python3", str(skill_root / "scripts" / "build_index.py")],
            check=True,
        )

    print(json.dumps({"ingested": len(records), "source": str(source)}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
